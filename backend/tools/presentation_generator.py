from typing import Dict, Any, List
from loguru import logger

class PresentationGenerator:
    async def generate_slides(self, topic: str, num_slides: int = 5) -> Dict[str, Any]:
        logger.info(f"Generating {num_slides} slides for: {topic}")
        try:
            from brain.model_router import ModelRouter
            router = ModelRouter()
            prompt = (
                f"Create an outline for a {num_slides}-slide presentation about '{topic}'. "
                "Return JSON array with objects having: title, bullet_points (array). "
                "No markdown, no explanations."
            )
            result = router.async_route_and_generate(prompt, task_type="general", max_cost=0.02)
            text = result.get("text", "") if isinstance(result, dict) else ""
            slides: List[Dict[str, Any]] = []
            import json
            try:
                cleaned = text.strip()
                if cleaned.startswith("```"):
                    cleaned = "\n".join(cleaned.splitlines()[1:])
                if cleaned.endswith("```"):
                    cleaned = "\n".join(cleaned.splitlines()[:-1])
                slides = json.loads(cleaned)
                if not isinstance(slides, list):
                    slides = []
            except Exception:
                for i in range(1, num_slides + 1):
                    slides.append({
                        "slide_number": i,
                        "title": f"Key Point {i} for {topic}",
                        "bullet_points": ["Detail A", "Detail B", "Detail C"],
                        "image_placeholder": "Relevant generated image",
                    })
            file_url = ""
            try:
                from pptx import Presentation as PptxPresentation
                prs = PptxPresentation()
                for slide_data in slides:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = slide_data.get("title", "")
                    body = slide.placeholders[1]
                    body.text = "\n".join(slide_data.get("bullet_points", []))
                out_path = f"data/{hash(topic)}.pptx"
                prs.save(out_path)
                file_url = out_path
            except Exception as pptx_err:
                logger.warning(f"PPTX generation failed: {pptx_err}")
                file_url = f"https://cdn.supremeai.example/presentations/{hash(topic)}.pptx"
            return {
                "status": "success",
                "topic": topic,
                "total_slides": len(slides),
                "slides_preview": slides,
                "download_url": file_url,
            }
        except Exception as exc:
            logger.error(f"Presentation generation failed: {exc}")
            return {
                "status": "error",
                "topic": topic,
                "error": str(exc),
            }
